import os
import re
from pathlib import Path


def clean_text(text: str) -> str:
    """Nettoie le texte extrait : supprime les espaces multiples et lignes vides."""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    return text.strip()


def extract_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return clean_text(f.read())


def extract_from_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return clean_text(text)


def extract_from_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return clean_text(text)


def extract_from_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    text = ""
    for sheet in wb.worksheets:
        text += f"[Feuille : {sheet.title}]\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text += row_text + "\n"
    return clean_text(text)


def extract_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"[Slide {i + 1}]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return clean_text(text)


def extract_from_image(file_path: str) -> str:
    """
    Extraction depuis image.
    Stratégie 1 : OCR avec pytesseract si disponible.
    Stratégie 2 : retourne un placeholder pour traitement multimodal.
    """
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang="fra+eng")
        return clean_text(text)
    except ImportError:
        # pytesseract non installé — signale que l'image doit être traitée par un modèle multimodal
        return f"[IMAGE:{file_path}]"
    except Exception as e:
        return f"[IMAGE_ERROR:{str(e)}]"


# Mapping extension -> fonction d'extraction
EXTRACTORS = {
    ".txt": extract_from_txt,
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".doc": extract_from_docx,
    ".xlsx": extract_from_xlsx,
    ".xls": extract_from_xlsx,
    ".pptx": extract_from_pptx,
    ".ppt": extract_from_pptx,
    ".png": extract_from_image,
    ".jpg": extract_from_image,
    ".jpeg": extract_from_image,
    ".webp": extract_from_image,
}


def extract(file_path: str) -> dict:
    """
    Point d'entrée principal.
    Retourne un dict avec le texte extrait, le type détecté et les métadonnées.
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    if not path.exists():
        return {
            "success": False,
            "error": f"Fichier non trouvé : {file_path}",
            "text": "",
            "file_type": extension,
            "requires_multimodal": False,
        }

    if extension not in EXTRACTORS:
        return {
            "success": False,
            "error": f"Type de fichier non supporté : {extension}",
            "text": "",
            "file_type": extension,
            "requires_multimodal": False,
        }

    try:
        text = EXTRACTORS[extension](file_path)

        # Détecter si l'image nécessite un modèle multimodal
        requires_multimodal = text.startswith("[IMAGE:")

        return {
            "success": True,
            "error": None,
            "text": text,
            "file_type": extension,
            "requires_multimodal": requires_multimodal,
            "char_count": len(text),
        }
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "text": "",
            "file_type": extension,
            "requires_multimodal": False,
        }


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python3 extractor.py <fichier>")
        sys.exit(1)
    result = extract(sys.argv[1])
    print(f"Type     : {result['file_type']}")
    print(f"Succès   : {result['success']}")
    if result['success']:
        print(f"Caractères : {result['char_count']}")
        print(f"Multimodal : {result['requires_multimodal']}")
        print(f"\nExtrait (500 premiers caractères) :\n{result['text'][:500]}")
    else:
        print(f"Erreur   : {result['error']}")